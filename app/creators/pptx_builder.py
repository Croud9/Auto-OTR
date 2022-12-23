import os, datetime
import pandas as pd
import numpy as np
import aspose.slides as slides
from PIL import Image
from pptx import Presentation
from os.path import isfile, join
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from transliterate import translit

class PPTXBulider():
    def __init__(self):
        self.input_data()
    
    def input_data(self):
        self.measure_units = {'kV': 'кВ', 'W': 'Вт', 'kW': 'кВт', 'MW': 'МВт', 'GW': 'ГВт', 'kWp': 'кВтп',
                    'year': 'год', 'GWh': 'ГВт*ч', 'MWh': 'МВт*ч', 'kWh': 'кВт*ч'}
        self.locale_ru = {'header_1': 'ТЕХНИЧЕСКИЕ ПАРАМЕТРЫ СОЛНЕЧНОЙ ЭЛЕКТРОСТАНЦИИ ПО ПРОЕКТУ ',
                        'header_2': 'ПОКАЗАТЕЛИ ВЫРАБОТКИ СЭС', 'header_3': 'ПРИМЕР ПРИНЦИПИАЛЬНОЙ СХЕМЫ', 
                        'months': {1: 'Январь', 2: 'Февраль', 3: 'Март', 4: 'Апрель', 5: 'Май',
                                    6: 'Июнь', 7: 'Июль', 8: 'Август', 9: 'Сентябрь', 10: 'Октябрь',
                                    11: 'Ноябрь', 12: 'Декабрь'},
                        'sides_world': {0: 'Север', 360: 'Север', 180: 'Юг', 270: 'Запад', 90: 'Восток',
                                        '0-90': 'Северо-Восток', '90-180': 'Юго-Восток', '180-270': 'Юго-Запад',
                                        '270-360': 'Северо-Запад'},
                        'title_column_chart': 'Помесячная выработка СЭС ', 'title_line_chart': 'Дневная выработка СЭС ',
                        'title_summer_series': 'Летний день', 'title_winter_series': 'Зимний день', 'title_spring_series': 'Весенний день',
                        'not_found': 'Н/Д', 'title_produced_energy': 'Годовая выработка электроэнергии', 'title_specific_production': 'Удельная годовая выработка',
                        'titles_table': ['ПАРАМЕТРЫ СЭС', 'Установленная мощность СЭС', 'Тип опорной конструкции', 'Ориентация и угол наклона',
                                        'Требуемая площадь', 'Уровень напряжения / Частота', 'Количество ФЭМ', 'Тип ФЭМ', 'Единичная мощность ФЭМ'],
                        'pcs': 'шт.', 'm2': 'м2', 'hz_and_kw': 'кВ / 50 Гц', 'w': 'Вт'}
        self.locale_en = {'header_1': 'SOLAR POWER PLANT (SPP) TECHNICAL PARAMETERS',
                        'header_2': 'ENERGY YIELD', 'header_3': 'EXAMPLE OF SINGLE LINE DIAGRAMM', 
                        'months': {1: 'January', 2: 'February', 3: 'March', 4: 'April', 5: 'May',
                                    6: 'June', 7: 'July', 8: 'August', 9: 'September', 10: 'October',
                                    11: 'November', 12: 'December'},
                        'sides_world': {0: 'North', 360: 'North', 180: 'South', 270: 'West', 90: 'East',
                                        '0-90': 'Northeast', '90-180': 'Southeast', '180-270': 'Southwest',
                                        '270-360': 'Northwest'},
                        'title_column_chart': 'Energy eield per month', 'title_line_chart': 'Energy eield per day',
                        'title_summer_series': 'Summer day', 'title_winter_series': 'Winter day', 'title_spring_series': 'Spring day',
                        'not_found': 'N/A', 'title_produced_energy': 'Annual energy yield', 'title_specific_production': 'Specific yield',
                        'titles_table': ['SES PARAMETERS', 'Installed capacity of SPP', 'Type of mounting structure', 'Oientation and tilt',
                                        'Required area', 'Voltage level / Frequency', 'Number of PV panel', 'PV type', 'Peak power of PV panel'],
                        'pcs': 'pcs', 'm2': 'm2', 'hz_and_kw': 'kW / 50 Hz', 'w': 'W'}

    def side_by_azimuth(self, azimuth):
        if azimuth < 0: azimuth += 360
        if azimuth in self.locale['sides_world']: return self.locale['sides_world'][azimuth]
        elif azimuth in range(0, 90): return self.locale['sides_world']['0-90']
        elif azimuth in range(270, 360): return self.locale['sides_world']['270-360']
        elif azimuth in range(90, 180): return self.locale['sides_world']['90-180']
        elif azimuth in range(180, 270): return self.locale['sides_world']['180-270']

    def daily_data(self, path):
        daily_measure_unit = self.locale['not_found']
        try:
            with open(path, 'r') as fp:
                for l_no, line in enumerate(fp):
                    if '[' in line: 
                        daily_measure_unit = line.split('[')[1].split(']')[0]

            if self.language == 'RU' and daily_measure_unit in self.measure_units.keys():
                daily_measure_unit = self.measure_units[daily_measure_unit]
            
            data_daily = pd.read_csv(path, skiprows = 3, header = None)
            df = pd.DataFrame(data_daily)

            april_data = list(df.iloc[np.flatnonzero(df[[0]] == 'April')].to_numpy()[0])
            july_data = list(df.iloc[np.flatnonzero(df[[0]] == 'July')].to_numpy()[0])
            december_data = list(df.iloc[np.flatnonzero(df[[0]] == 'December')].to_numpy()[0])
            del april_data[-1]
            del july_data[-1]
            del december_data[-1]

            april_num = [0 if i < 0 else i for i in april_data[1: ]]
            july_num = [0 if i < 0 else i for i in july_data[1: ]]
            december_num = [0 if i < 0 else i for i in december_data[1: ]]
            return {'april': april_num, 'july': july_num, 'december': december_num, 'measure_unit': daily_measure_unit, 'broken_file': False}
        except Exception:
            return {'broken_file': True}

    def del_slide(self, currencies_ppt, index_slide):
        xml_slides = currencies_ppt.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[index_slide]) 

    def chart_config(self, chart, measure_unit, type_chart = 'column'):
        category_axis = chart.category_axis #Горизонтальная ось
        value_axis = chart.value_axis #Вертикальная ось
        category_axis.position = XL_LABEL_POSITION.OUTSIDE_END
        if type_chart == 'column': 
            category_axis.tick_labels.font.size = Pt(15)
            value_axis.axis_title.text_frame.text = measure_unit
        elif type_chart == 'line':
            category_axis.tick_labels.font.size = Pt(11)
            value_axis.axis_title.text_frame.text = measure_unit

        value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(15)
        value_axis.axis_title.text_frame.paragraphs[0].font.bold = False
        value_axis.tick_labels.font.size = Pt(14)

    def small_chart(self, x_pos, y_pos, type_chart, data, slide, p, measure_unit, count_schemes_on_slide):
        if count_schemes_on_slide == 4:
            cx, cy = Inches(7.5), Inches(3.75)
        elif count_schemes_on_slide == 2:
            cx, cy = Inches(8.5), Inches(7)
        chart_data = ChartData()

        if type_chart == 'column': 
            chart_data.categories = self.locale['months'].values()
            chart_data.add_series('1', data)
            xl_type_chart = XL_CHART_TYPE.COLUMN_CLUSTERED 
        elif type_chart == 'line':
            chart_data.categories = range(0, 24)
            chart_data.add_series(self.locale['title_summer_series'], data['july'])
            chart_data.add_series(self.locale['title_winter_series'], data['december'])
            chart_data.add_series(self.locale['title_spring_series'], data['april'])
            xl_type_chart = XL_CHART_TYPE.LINE
        
        graphic_frame = slide.shapes.add_chart(xl_type_chart, Inches(x_pos), Inches(y_pos), cx, cy, chart_data)
        chart = graphic_frame.chart
        chart.has_legend = False

        if type_chart == 'column':
            chart.plots[0].vary_by_categories = False
            title = f"{self.locale['title_column_chart']} {p}"
        elif type_chart == 'line':
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = Pt(14)
            chart.legend.include_in_layout = False
            chart.series[0].smooth = True
            chart.series[1].smooth = True
            chart.series[2].smooth = True
            title = f"{self.locale['title_line_chart']} {p}"
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(15)
        chart.chart_title.text_frame.paragraphs[0].font.bold = False  
        self.chart_config(chart, measure_unit, type_chart)

    def footer(self, slide):
        shapes = slide.shapes
        text_box_list = []
        for shape_idx in range(len(shapes)):
            shape = shapes[shape_idx]
            if shape.shape_type == 17:
                text_box_list.append(shape_idx)

        date = datetime.date.today()
        footer = shapes[text_box_list[0]].text_frame.paragraphs[0]
        footer.runs[0].text = f"{self.title_project} / {self.locale['months'][date.month]} {date.year}"

    def header(self, slide, text):
        shapes = slide.shapes
        placeholder_list = []

        for shape_idx in range(len(shapes)):
            shape = shapes[shape_idx]
            if shape.shape_type == 14:
                placeholder_list.append(shape_idx)
                
        header = shapes[placeholder_list[0]].text_frame.paragraphs[0]
        header.runs[0].text = text

    def main_table(self, data, slide, count_pv, keys_pvs):
        if type(data['pvsyst']['produced_energy']) is dict:
            produced_energy = []
            specific_production = []
            row = 13
            produced_energy.append(f"{data['pvsyst']['produced_energy']['(P50)']} {self.prod_e_ru} (P50)")
            produced_energy.append(f"{data['pvsyst']['produced_energy'][data['current_p']]} {self.prod_e_ru} {data['current_p']}")
            specific_production.append(f"{data['pvsyst']['specific_production']['(P50)']} {self.spec_p_ru} (P50)")
            specific_production.append(f"{data['pvsyst']['specific_production'][data['current_p']]} {self.spec_p_ru} {data['current_p']}")
            last_rows = [self.locale['title_produced_energy'], '', self.locale['title_specific_production'], ''] 
        else:
            row = 11
            last_rows = [self.locale['title_produced_energy'], self.locale['title_specific_production']]

        x, y, cx, cy = Inches (1.2), Inches (2.2), Inches (16.8), Inches (5.8)
        columns = 1 + count_pv
        ses_table = slide.shapes.add_table(row, columns, x, y, cx, cy).table

        tilt_and_azimuth = ''
        num = 0
        for value in data['pvsyst']['tilts_and_azimuths'].values():
            num += 1
            if num > 1:
                tilt_and_azimuth += ' + '
            tilt_and_azimuth += f"{self.side_by_azimuth(value['azimuth'])}, {value['tilt']}°"

        self.locale['titles_table'].extend(last_rows)
        params = [f"{data['pvsyst']['pnom_PV']} {self.pnom_PV_measure_unit}", data['roof'], tilt_and_azimuth, f"{data['pvsyst']['module_area']} {self.locale['m2']}",
                    f"{data['pvsyst']['voltage_level']} {self.locale['hz_and_kw']}", f"{data['pvsyst']['nb_PV']} {self.locale['pcs']}"]

        for i in range(row): # row
            cell = ses_table.cell(i, 0)
            para = cell.text_frame.paragraphs[0]
            para.text = f"{self.locale['titles_table'][i]}"
            para.font.size = Pt(18)

        for i in range(6): # row
            cell = ses_table.cell(i + 1, 1)
            para = cell.text_frame.paragraphs[0]
            para.text = f"{params[i]}"
            para.font.size = Pt(18)
            para.alignment = PP_ALIGN.CENTER

        if count_pv > 1:
            for row_pv in range(2): # row
                col = 1
                for pv in keys_pvs: #column
                    cell = ses_table.cell(7 + row_pv, col)
                    para = cell.text_frame.paragraphs[0]
                    if row_pv == 0:
                        para.text = data['pvs'][pv]['type_pv']
                    else:
                        para.text = f"{data['pvs'][pv]['p_nom_pv']} {self.locale['w']}"
                    para.font.size = Pt(18)
                    para.alignment = PP_ALIGN.CENTER
                    col += 1

            for i in range(1, 7):
                start_cell = ses_table.rows[i].cells[1]	
                ofset_cell = ses_table.cell(i, -1)
                start_cell.merge(ofset_cell)

            for i in range(9, row):
                start_cell = ses_table.rows[i].cells[1]	
                ofset_cell = ses_table.cell(i, -1)
                start_cell.merge(ofset_cell)

            for i in range(1, count_pv + 1):
                size_column = 12.8 / count_pv
                ses_table.columns[i].width = Inches(size_column)
        else:
            ses_table.columns[1].width = Inches(5.4)
            for i in range(2): # row
                cell = ses_table.cell(7 + i, 1)
                para = cell.text_frame.paragraphs[0]
                if i == 0:
                    para.text = data['pvs']['found_pv_0']['type_pv']
                else:
                    para.text = f"{data['pvs']['found_pv_0']['p_nom_pv']} {self.locale['w']}"
                para.font.size = Pt(18)
                para.alignment = PP_ALIGN.CENTER

        start_cell = ses_table.rows[0].cells[0]	
        ofset_cell = ses_table.cell(0, -1)
        start_cell.merge(ofset_cell)

        start_cell.fill.solid()
        start_cell.fill.fore_color.rgb = RGBColor(2, 82, 56) #025238
        start_cell.margin_left = Inches(0.2)
        start_cell.margin_top = Inches(0.2)

        for i in range(1, row):
            for j in range(count_pv + 1):
                if i % 2 == 0:
                    color = RGBColor(238, 239, 241) #eeeff1
                else:
                    color = RGBColor(255, 255, 255) #fff
                start_cell = ses_table.rows[i].cells[j]	
                start_cell.fill.solid()
                start_cell.fill.fore_color.rgb = color

        if type(data['pvsyst']['produced_energy']) is dict:
            for i in range(2):
                cell_prod = ses_table.cell(9 + i, 1)
                para = cell_prod.text_frame.paragraphs[0]
                para.text = f"{produced_energy[i]}"
                para.font.size = Pt(18)
                para.alignment = PP_ALIGN.CENTER

                cell_spec = ses_table.cell(11 + i, 1)
                para = cell_spec.text_frame.paragraphs[0]
                para.text = f"{specific_production[i]}"
                para.font.size = Pt(18)
                para.alignment = PP_ALIGN.CENTER

            start_cell = ses_table.rows[9].cells[0]	
            ofset_cell = ses_table.cell(10, 0)
            start_cell.merge(ofset_cell)

            start_cell = ses_table.rows[11].cells[0]	
            ofset_cell = ses_table.cell(12, 0)
            start_cell.merge(ofset_cell)
        else:
            cell_prod = ses_table.cell(9, 1)
            para = cell_prod.text_frame.paragraphs[0]
            para.text = f"{data['pvsyst']['produced_energy']} {self.prod_e_ru}"
            para.font.size = Pt(18)
            para.alignment = PP_ALIGN.CENTER

            cell_spec = ses_table.cell(10, 1)
            para = cell_spec.text_frame.paragraphs[0]
            para.text = f"{data['pvsyst']['specific_production']} {self.spec_p_ru}"
            para.font.size = Pt(18)
            para.alignment = PP_ALIGN.CENTER

        ses_table.columns[0].width = Inches(4)
        for i in range(1, row):
            ses_table.rows[i].height = Inches(0.45)
        ses_table.rows[0].height = Inches(0.8)
        ses_table.rows[3].height = Inches(0.8)
        if row == 11:
            ses_table.rows[9].height = Inches(0.8)
            height_table = (row - 3) * 0.45 + 3 * 0.8
        else:
            height_table = (row - 2) * 0.45 + 2 * 0.8
        return height_table

    def slide_1(self, data, currencies_ppt):
        slide = currencies_ppt.slides[0]
        count_pv = len(data['pvs'])
        keys_pvs = []
        for key in data['pvs'].keys():
            keys_pvs.append(key)
        height_table = self.main_table(data, slide, count_pv, keys_pvs)
        if count_pv <= 1:
            img_path = 'Data\System\PPTX\img_pvs.jpg'
            slide.shapes.add_picture(img_path, Inches (11.1), Inches (2.2), width = Inches(6.9), height = Inches(height_table))
        self.header(slide, f"{self.locale['header_1']} {self.title_project}")
        self.footer(slide)

    def slide_2(self, data, currencies_ppt):
        slide = currencies_ppt.slides[1]
        x, y, cx, cy = Inches(2), Inches(2), Inches(14), Inches(6)
        chart_data = ChartData()
        chart_data.categories = self.locale['months'].values()
        chart_data.add_series('1', data['pvsyst']['months_e_grid'])
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        graphic_frame.left = (currencies_ppt.slide_width - graphic_frame.width) // 2
        chart = graphic_frame.chart
        chart.chart_title.text_frame.text = self.locale['title_column_chart']
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
        chart.chart_title.text_frame.paragraphs[0].font.bold = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.vary_by_categories = False
        data_labels = plot.data_labels #Значения столбцов
        data_labels.font.size = Pt(13)
        data_labels.font.bold = True
        data_labels.font.color.rgb = RGBColor(89, 89, 89)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        self.chart_config(chart, self.e_grid_measure_unit)
        self.header(slide, self.locale['header_2'])
        self.footer(slide)

    def slide_2_four_schemes(self, data, currencies_ppt, data_day):
        slide = currencies_ppt.slides[1]
        count_schemes_on_slide = 4
        self.small_chart(1.25, 1.75, 'column', data['pvsyst']['months_e_grid'], slide, '(P50)', self.e_grid_measure_unit, count_schemes_on_slide)
        self.small_chart(9.75, 1.75, 'line', data_day, slide, '(P50)', data_day['measure_unit'], count_schemes_on_slide)
        p_50 = float(data['pvsyst']['specific_production']['(P50)'])
        p_select = float(data['pvsyst']['specific_production'][data['current_p']])
        data_months_p = list(map(lambda x: x * (p_select / p_50), data['pvsyst']['months_e_grid']))
        data_day_p = {'july': list(map(lambda x: x * (p_select / p_50), data_day['july'])),
                        'december': list(map(lambda x: x * (p_select / p_50), data_day['december'])),
                        'april': list(map(lambda x: x * (p_select / p_50), data_day['april']))}
        self.small_chart(1.25, 5.75, 'column', data_months_p, slide, data['current_p'], self.e_grid_measure_unit, count_schemes_on_slide)
        self.small_chart(9.75, 5.75, 'line', data_day_p, slide, data['current_p'], data_day['measure_unit'], count_schemes_on_slide)
        self.header(slide, self.locale['header_2'])
        self.footer(slide)

    def slide_2_two_schemes(self, data, currencies_ppt, data_day):
        slide = currencies_ppt.slides[1]
        count_schemes_on_slide = 2
        self.small_chart(1, 2, 'column', data['pvsyst']['months_e_grid'], slide, '(P50)', self.e_grid_measure_unit, count_schemes_on_slide)
        self.small_chart(9.5, 2, 'line', data_day, slide, '(P50)', data_day['measure_unit'], count_schemes_on_slide)
        self.header(slide, self.locale['header_2'])
        self.footer(slide)

    def slide_3(self, currencies_ppt, index_slide, img_path):
        slide = currencies_ppt.slides[index_slide]

        img = Image.open(img_path)
        h = img.height
        if 1900 < h <= 2400:
            height_in_slide = Inches(6)
        elif h <= 1900:
            height_in_slide = Inches(5)
        else:
            height_in_slide = Inches(7.9)

        top_pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1.75), height = height_in_slide)
        top_pic.left = (currencies_ppt.slide_width - top_pic.width) // 2
        self.header(slide, self.locale['header_3'])
        self.footer(slide)

    def copy_slide(self, count_slide):
        with slides.Presentation("Data\System\PPTX\layout.pptx") as pres1:
            with slides.Presentation("Data\System\PPTX\layout.pptx") as pres2:
                for i in range(count_slide - 1):
                    pres1.slides.add_clone(pres2.slides[2])
                pres1.save("Data\System\PPTX\layout_merge.pptx", slides.export.SaveFormat.PPTX)

        ppt = Presentation("Data\System\PPTX\layout_merge.pptx")
        for slide in ppt.slides:
            shapes = slide.shapes
            for shape in shapes:
                if shape.shape_type == 17 and shape.height == 1569660:
                    shapes.element.remove(shape.element)
        ppt.save("Data\System\PPTX\layout_merge.pptx")

    def translation(self, data):
        self.prod_e_ru = ''
        self.spec_p_ru = ''
        prod_e = data['pvsyst']['produced_and_specific_measure_units'][0].split('/')
        spec_p = data['pvsyst']['produced_and_specific_measure_units'][1].split('/')
        for i in prod_e:
            if i in self.measure_units.keys():
                self.prod_e_ru += self.measure_units[i]
            else:
                self.prod_e_ru += i
            if i != prod_e[-1]:
                self.prod_e_ru += '/'

        for i in spec_p:
            if i in self.measure_units.keys():
                self.spec_p_ru += self.measure_units[i]
            else:
                self.spec_p_ru += i
            if i != prod_e[-1]:
                self.spec_p_ru += '/'

        if data['pvsyst']['pnom_PV_measure_unit'] in self.measure_units.keys():
            self.pnom_PV_measure_unit = self.measure_units[data['pvsyst']['pnom_PV_measure_unit']]
        else:
            self.pnom_PV_measure_unit= data['pvsyst']['pnom_PV_measure_unit']

        if data['pvsyst']['e_grid_measure_unit'] in self.measure_units.keys():
            self.e_grid_measure_unit = self.measure_units[data['pvsyst']['e_grid_measure_unit']]
        else:
            self.e_grid_measure_unit = data['pvsyst']['e_grid_measure_unit'] 

    def localization(self, data):
        if self.language == 'RU':
            fp_folder = 'Data/System/Images/PPTX/RU'
            self.locale = self.locale_ru
            self.title_project = data['title_project']
            self.translation(data)
        else:
            fp_folder = 'Data/System/Images/PPTX/EN'
            self.locale = self.locale_en
            self.title_project = translit(data['title_project'], "ru", reversed=True).upper()
            self.prod_e_ru = data['pvsyst']['produced_and_specific_measure_units'][0]
            self.spec_p_ru = data['pvsyst']['produced_and_specific_measure_units'][1]
            self.pnom_PV_measure_unit= data['pvsyst']['pnom_PV_measure_unit']
            self.e_grid_measure_unit = data['pvsyst']['e_grid_measure_unit'] 
        return fp_folder

    def create_pptx(self, data):
        status = 'ok'
        self.language = data['locale']
        fp_folder = self.localization(data)
        files = [f for f in os.listdir(fp_folder) if isfile(join(fp_folder, f))]
        if len(files) != 0:
            if len(files) > 1:
                self.copy_slide(len(files))
                currencies_ppt = Presentation("Data\System\PPTX\layout_merge.pptx")
            else:
                currencies_ppt = Presentation("Data\System\PPTX\layout.pptx")

            index_scheme_slide = 2
            for file in files:
                img_path = fp_folder + f"/{file}"
                self.slide_3(currencies_ppt, index_scheme_slide, img_path)
                index_scheme_slide += 1

            self.slide_1(data, currencies_ppt)
            if type(data['pvsyst']['specific_production']) is dict and data['type_schemes'] == 2:
                data_day = self.daily_data(data['path_to_daily_csv'])
                if data_day['broken_file'] != True:
                    self.slide_2_four_schemes(data, currencies_ppt, data_day)
                else:
                    status = 'fail'
            elif data['type_schemes'] == 1:
                data_day = self.daily_data(data['path_to_daily_csv'])
                if data_day['broken_file'] != True:
                    self.slide_2_two_schemes(data, currencies_ppt, data_day)
                else:
                    status = 'fail'
            else:
                self.slide_2(data, currencies_ppt)

            currencies_ppt.save("Data\Result\Summary.pptx")
        return status