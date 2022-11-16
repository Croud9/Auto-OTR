# from pptx import Presentation
# from pptx.util import Inches
# import win32com.client
# import os

# currencies_ppt = Presentation('Currencies.pptx')
# slide = currencies_ppt.slides[0]
# shapes = slide.shapes
# for shape in shapes:
#     if shape.shape_type == 13:
#         shapes.element.remove(shape.element)
# top_img_path = 'top.png'
# bottom_img_path = 'bottom.png'
# top_pic = slide.shapes.add_picture(top_img_path, Inches(0.40), Inches(4.85), width=Inches(5.30))
# bottom_pic = slide.shapes.add_picture(bottom_img_path, Inches(5.25), Inches(4.85), width=Inches(5.30))
# ref_element = slide.shapes[0]._element
# ref_element.addprevious(top_pic._element)
# ref_element.addprevious(bottom_pic._element)
# shapes = slide.shapes
# text_box_list = []
# auto_shape_list = []
# table_list = []

# for shape_idx in range(len(shapes)):
#     shape = shapes[shape_idx]
#     if shape.shape_type == 17:
#         text_box_list.append(shape_idx)
#     if shape.shape_type == 1:
#         auto_shape_list.append(shape_idx)
#     if shape.shape_type == 19:
#         table_list.append(shape_idx)

# last_update_date_textbox_height = max([shapes[shape_idx].height for shape_idx in text_box_list])
# last_update_date_idx = [shape_idx for shape_idx in text_box_list if shapes[shape_idx].height == last_update_date_textbox_height][0]
# top_label_left = min([shapes[shape_idx].left for shape_idx in auto_shape_list])
# top_label_idx = [shape_idx for shape_idx in auto_shape_list if shapes[shape_idx].left == top_label_left][0]
# auto_shape_list.remove(top_label_idx)
# bottom_label_idx = auto_shape_list[0]
# top_table_left = min([shapes[shape_idx].left for shape_idx in table_list])
# top_table_idx = [shape_idx for shape_idx in table_list if shapes[shape_idx].left == top_table_left][0]
# table_list.remove(top_table_idx)
# bottom_table_idx = table_list[0]
# paragraph = shapes[last_update_date_idx].text_frame.paragraphs[0]
# paragraph.runs[4].text = datetime_now.strftime("%#d %b %Y %H:%M")
# paragraph = shapes[top_label_idx].text_frame.paragraphs[0]
# paragraph.runs[0].text = top_df['Name'][0].replace('/', ' / ')
# paragraph = shapes[bottom_label_idx].text_frame.paragraphs[0]
# paragraph.runs[0].text = bottom_df['Name'][0].replace('/', ' / ')
# top_table = shapes[top_table_idx].table

# for i in range(5):
#     for j in range(4):
#         cell = top_table.cell(i+1, j)
#         paragraph = cell.text_frame.paragraphs[0]
#         run = paragraph.runs[0]
#         run.text = str(top_df.iloc[i, j])
# bottom_table = shapes[bottom_table_idx].table

# for i in range(5):
#     for j in range(4):
#         cell = bottom_table.cell(i+1, j)
#         paragraph = cell.text_frame.paragraphs[0]
#         run = paragraph.runs[0]
#         run.text = str(bottom_df.iloc[i, j])
# currencies_ppt.save('New_Currencies.pptx')















from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

currencies_ppt = Presentation("Data\GEN PPTX\Layout_for_summary.pptx")
# prs.slide_width = Inches(16)
# prs.slide_height = Inches(9)
slide = currencies_ppt.slides[0]
title = slide.shapes.title # assigning a title

for shape in slide.shapes:
    print(shape.shape_type)

# shapes = slide.shapes
# for shape in shapes:
#     if shape.shape_type == 13:
#         shapes.element.remove(shape.element)
currencies_ppt.save("Data\GEN PPTX\Summary.pptx") # saving file


